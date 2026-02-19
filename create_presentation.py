"""
Auto Insurance Portfolio Analysis - Presentation Generator
Creates a professional PPTX to demo AI agent capabilities.
"""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
BG_DARK = RGBColor(0x0F, 0x11, 0x17)
BG_SURFACE = RGBColor(0x1A, 0x1D, 0x27)
BG_SURFACE2 = RGBColor(0x23, 0x27, 0x33)
BORDER = RGBColor(0x2E, 0x33, 0x45)
TEXT_WHITE = RGBColor(0xE4, 0xE6, 0xED)
TEXT_MUTED = RGBColor(0x8B, 0x8F, 0xA3)
ACCENT = RGBColor(0x63, 0x66, 0xF1)
ACCENT_LIGHT = RGBColor(0x81, 0x8C, 0xF8)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
PINK = RGBColor(0xEC, 0x48, 0x99)
ORANGE = RGBColor(0xF9, 0x73, 0x16)

# ── Load data ──
with open("dashboard_data.json") as f:
    raw = json.load(f)

# Parse: [id, gender(1=M/0=F), age, car_year, premium, loss]
data = []
for r in raw:
    data.append({
        "id": r[0], "gender": "Male" if r[1] == 1 else "Female",
        "age": r[2], "car_year": r[3], "premium": r[4], "loss": r[5]
    })

def age_group(a):
    if a < 25: return "18-24"
    if a < 35: return "25-34"
    if a < 45: return "35-44"
    if a < 55: return "45-54"
    if a < 65: return "55-64"
    return "65+"

def car_era(y):
    if y < 2005: return "2000-04"
    if y < 2010: return "2005-09"
    if y < 2015: return "2010-14"
    if y < 2020: return "2015-19"
    return "2020-25"

AGE_GROUPS = ["18-24", "25-34", "35-44", "45-54", "55-64", "65+"]
CAR_ERAS = ["2000-04", "2005-09", "2010-14", "2015-19", "2020-25"]

# ── Aggregations ──
total_premium = sum(d["premium"] for d in data)
total_loss = sum(d["loss"] for d in data)
claims_count = sum(1 for d in data if d["loss"] > 0)
loss_ratio = total_loss / total_premium
claims_freq = claims_count / len(data)

by_age = {}
for ag in AGE_GROUPS:
    rows = [d for d in data if age_group(d["age"]) == ag]
    prem = sum(r["premium"] for r in rows)
    loss = sum(r["loss"] for r in rows)
    claims = sum(1 for r in rows if r["loss"] > 0)
    by_age[ag] = {"count": len(rows), "premium": prem, "loss": loss, "claims": claims,
                  "avg_premium": prem / len(rows) if rows else 0,
                  "loss_ratio": loss / prem if prem else 0,
                  "claims_freq": claims / len(rows) if rows else 0}

by_era = {}
for era in CAR_ERAS:
    rows = [d for d in data if car_era(d["car_year"]) == era]
    prem = sum(r["premium"] for r in rows)
    loss = sum(r["loss"] for r in rows)
    by_era[era] = {"count": len(rows), "loss_ratio": loss / prem if prem else 0}

by_gender = {}
for g in ["Male", "Female"]:
    rows = [d for d in data if d["gender"] == g]
    prem = sum(r["premium"] for r in rows)
    loss = sum(r["loss"] for r in rows)
    claims = sum(1 for r in rows if r["loss"] > 0)
    by_gender[g] = {"count": len(rows), "premium": prem, "loss": loss,
                    "avg_premium": prem / len(rows),
                    "loss_ratio": loss / prem if prem else 0,
                    "claims_freq": claims / len(rows)}

# Heatmap data
heatmap = {}
for ag in AGE_GROUPS:
    heatmap[ag] = {}
    for era in CAR_ERAS:
        rows = [d for d in data if age_group(d["age"]) == ag and car_era(d["car_year"]) == era]
        prem = sum(r["premium"] for r in rows)
        loss = sum(r["loss"] for r in rows)
        heatmap[ag][era] = {"n": len(rows), "lr": loss / prem if prem else 0}

# ── Presentation ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ──
def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox

def add_multiline(slide, left, top, width, height, lines, font_size=14,
                  color=TEXT_WHITE, line_spacing=1.5, alignment=PP_ALIGN.LEFT):
    """lines: list of (text, bold, color) tuples"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, (text, bold, c) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = c
        p.font.bold = bold
        p.font.name = "Arial"
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))

# ── SLIDE 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Accent bar at top
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)

# Title text
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Auto Insurance Portfolio Analysis", font_size=44, color=ACCENT_LIGHT, bold=True,
             alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1), Inches(3.1), Inches(11), Inches(0.6),
             "AI-Powered Risk Insights  |  10,000 Policyholders", font_size=20, color=TEXT_MUTED,
             alignment=PP_ALIGN.CENTER)

# Divider line
add_rect(slide, Inches(4.5), Inches(3.9), Inches(4.3), Inches(0.02), BORDER)

# Bottom info
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "Generated by Claude Code  |  February 2026", font_size=14, color=TEXT_MUTED,
             alignment=PP_ALIGN.CENTER)

# Bottom accent bar
add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), CYAN)


# ── SLIDE 2: Executive Summary ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "Executive Summary", font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8), Inches(0.4),
             "Portfolio overview across 10,000 active policies", font_size=14, color=TEXT_MUTED)

# KPI Cards
kpis = [
    ("10,000", "Active Policies", ACCENT_LIGHT),
    (f"${total_premium/1e6:.1f}M", "Total Premium", CYAN),
    (f"${total_loss/1e6:.1f}M", "Total Claims Paid", AMBER),
    (f"{loss_ratio*100:.1f}%", "Loss Ratio", RED if loss_ratio > 0.8 else AMBER),
    (f"{claims_freq*100:.1f}%", "Claims Frequency", GREEN),
]

card_w = Inches(2.2)
card_h = Inches(1.6)
start_x = Inches(0.8)
gap = Inches(0.27)

for i, (value, label, color) in enumerate(kpis):
    x = start_x + i * (card_w + gap)
    y = Inches(1.7)
    card = add_shape(slide, x, y, card_w, card_h, BG_SURFACE, BORDER)

    # Accent line at top of card
    add_rect(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.04), color)

    add_text_box(slide, x + Inches(0.15), y + Inches(0.35), card_w - Inches(0.3), Inches(0.3),
                 label, font_size=10, color=TEXT_MUTED)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.7), card_w - Inches(0.3), Inches(0.6),
                 value, font_size=28, color=color, bold=True)

# Key insight box
add_shape(slide, Inches(0.8), Inches(3.7), Inches(11.7), Inches(3.3), BG_SURFACE, BORDER)
add_text_box(slide, Inches(1.1), Inches(3.9), Inches(5), Inches(0.4),
             "Key Portfolio Metrics", font_size=16, color=TEXT_WHITE, bold=True)

insights = [
    ("Highest Risk Segment:", f" Young drivers (18-24) have {by_age['18-24']['claims_freq']*100:.0f}% claims frequency vs {by_age['35-44']['claims_freq']*100:.0f}% for 35-44 age group", AMBER),
    ("Gender Split:", f" {by_gender['Male']['count']:,} male ({by_gender['Male']['loss_ratio']*100:.0f}% LR) vs {by_gender['Female']['count']:,} female ({by_gender['Female']['loss_ratio']*100:.0f}% LR)", CYAN),
    ("Car Age Impact:", f" Older cars (2000-04) show {by_era['2000-04']['loss_ratio']*100:.0f}% loss ratio vs {by_era['2020-25']['loss_ratio']*100:.0f}% for newer (2020-25)", GREEN),
    ("Portfolio Health:", f" Overall 70.4% loss ratio is within industry acceptable range (60-80%)", ACCENT_LIGHT),
]

for i, (label, detail, color) in enumerate(insights):
    y_pos = Inches(4.5) + i * Inches(0.65)
    # Bullet dot
    add_rect(slide, Inches(1.1), y_pos + Inches(0.08), Inches(0.12), Inches(0.12), color)
    txbox = slide.shapes.add_textbox(Inches(1.4), y_pos, Inches(11), Inches(0.5))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = label
    run1.font.size = Pt(13)
    run1.font.color.rgb = color
    run1.font.bold = True
    run1.font.name = "Arial"
    run2 = p.add_run()
    run2.text = detail
    run2.font.size = Pt(13)
    run2.font.color.rgb = TEXT_WHITE
    run2.font.name = "Arial"


# ── SLIDE 3: Loss Ratio by Age Group (Chart) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), AMBER)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "Loss Ratio by Age Group", font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8), Inches(0.4),
             "Young and elderly drivers show highest loss ratios", font_size=14, color=TEXT_MUTED)

# Chart
chart_data = CategoryChartData()
chart_data.categories = AGE_GROUPS
chart_data.add_series("Loss Ratio", [by_age[ag]["loss_ratio"] * 100 for ag in AGE_GROUPS])

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.6),
    Inches(7.5), Inches(5.2), chart_data
)
chart = chart_frame.chart
chart.has_legend = False
chart.style = 2

# Style the chart
plot = chart.plots[0]
plot.gap_width = 80
series = plot.series[0]

# Color bars based on loss ratio
colors_map = []
for ag in AGE_GROUPS:
    lr = by_age[ag]["loss_ratio"]
    if lr > 0.8:
        colors_map.append(RED)
    elif lr > 0.7:
        colors_map.append(AMBER)
    else:
        colors_map.append(GREEN)

for i, c in enumerate(colors_map):
    pt = series.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = c

# Value axis formatting
val_axis = chart.value_axis
val_axis.has_title = True
val_axis.axis_title.text_frame.paragraphs[0].text = "Loss Ratio (%)"
val_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
val_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED
val_axis.minimum_scale = 0
val_axis.maximum_scale = 100

cat_axis = chart.category_axis
for lbl in [val_axis, cat_axis]:
    lbl.tick_labels.font.size = Pt(11)
    lbl.tick_labels.font.color.rgb = TEXT_MUTED

# Right side: insight cards
card_x = Inches(8.8)
card_w2 = Inches(4)

insights_right = [
    ("Highest Risk", f"{by_age['65+']['loss_ratio']*100:.0f}%", "65+ age group", RED),
    ("Lowest Risk", f"{by_age['25-34']['loss_ratio']*100:.0f}%", "25-34 age group", GREEN),
    ("Largest Pool", f"{by_age['35-44']['count']:,}", "35-44 policyholders", CYAN),
]

for i, (title, value, sub, color) in enumerate(insights_right):
    y = Inches(1.8) + i * Inches(1.8)
    add_shape(slide, card_x, y, card_w2, Inches(1.5), BG_SURFACE, BORDER)
    add_rect(slide, card_x, y, Inches(0.06), Inches(1.5), color)
    add_text_box(slide, card_x + Inches(0.3), y + Inches(0.15), Inches(3.5), Inches(0.3),
                 title, font_size=11, color=TEXT_MUTED)
    add_text_box(slide, card_x + Inches(0.3), y + Inches(0.5), Inches(3.5), Inches(0.5),
                 value, font_size=30, color=color, bold=True)
    add_text_box(slide, card_x + Inches(0.3), y + Inches(1.0), Inches(3.5), Inches(0.3),
                 sub, font_size=11, color=TEXT_MUTED)


# ── SLIDE 4: Claims Frequency + Premium (dual chart) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Claims Frequency & Average Premium", font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "Young drivers: highest claims rate AND highest premiums", font_size=14, color=TEXT_MUTED)

# Claims Frequency chart (left)
chart_data1 = CategoryChartData()
chart_data1.categories = AGE_GROUPS
chart_data1.add_series("Claims Freq %", [by_age[ag]["claims_freq"] * 100 for ag in AGE_GROUPS])

cf1 = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6),
    Inches(6), Inches(5.2), chart_data1
)
c1 = cf1.chart
c1.has_legend = False
c1.has_title = True
c1.chart_title.text_frame.paragraphs[0].text = "Claims Frequency (%)"
c1.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
c1.chart_title.text_frame.paragraphs[0].font.color.rgb = CYAN
c1.chart_title.text_frame.paragraphs[0].font.bold = True
plot1 = c1.plots[0]
plot1.gap_width = 80
for i in range(len(AGE_GROUPS)):
    plot1.series[0].points[i].format.fill.solid()
    plot1.series[0].points[i].format.fill.fore_color.rgb = CYAN
c1.value_axis.minimum_scale = 0
c1.value_axis.maximum_scale = 60
c1.value_axis.tick_labels.font.size = Pt(10)
c1.value_axis.tick_labels.font.color.rgb = TEXT_MUTED
c1.category_axis.tick_labels.font.size = Pt(10)
c1.category_axis.tick_labels.font.color.rgb = TEXT_MUTED

# Average Premium chart (right)
chart_data2 = CategoryChartData()
chart_data2.categories = AGE_GROUPS
chart_data2.add_series("Avg Premium $", [by_age[ag]["avg_premium"] for ag in AGE_GROUPS])

cf2 = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.8), Inches(1.6),
    Inches(6), Inches(5.2), chart_data2
)
c2 = cf2.chart
c2.has_legend = False
c2.has_title = True
c2.chart_title.text_frame.paragraphs[0].text = "Average Annual Premium ($)"
c2.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
c2.chart_title.text_frame.paragraphs[0].font.color.rgb = ACCENT_LIGHT
c2.chart_title.text_frame.paragraphs[0].font.bold = True
plot2 = c2.plots[0]
plot2.gap_width = 80
for i in range(len(AGE_GROUPS)):
    plot2.series[0].points[i].format.fill.solid()
    plot2.series[0].points[i].format.fill.fore_color.rgb = ACCENT_LIGHT
c2.value_axis.minimum_scale = 0
c2.value_axis.maximum_scale = 2500
c2.value_axis.tick_labels.font.size = Pt(10)
c2.value_axis.tick_labels.font.color.rgb = TEXT_MUTED
c2.category_axis.tick_labels.font.size = Pt(10)
c2.category_axis.tick_labels.font.color.rgb = TEXT_MUTED


# ── SLIDE 5: Risk Heatmap Table ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), RED)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Risk Heatmap: Age Group x Car Era", font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "Loss ratio matrix reveals risk concentration patterns", font_size=14, color=TEXT_MUTED)

# Build table
rows = len(AGE_GROUPS) + 1
cols = len(CAR_ERAS) + 1
table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
table = table_shape.table

# Column widths
table.columns[0].width = Inches(1.5)
for j in range(1, cols):
    table.columns[j].width = Inches(2.04)

# Header row
header_cell = table.cell(0, 0)
header_cell.text = "Age \\ Car Era"
for j, era in enumerate(CAR_ERAS):
    table.cell(0, j + 1).text = era

# Data rows
for i, ag in enumerate(AGE_GROUPS):
    table.cell(i + 1, 0).text = ag
    for j, era in enumerate(CAR_ERAS):
        lr = heatmap[ag][era]["lr"]
        n = heatmap[ag][era]["n"]
        cell = table.cell(i + 1, j + 1)
        cell.text = f"{lr*100:.0f}%  (n={n})"

# Style cells
def lr_color(lr):
    if lr < 0.5: return RGBColor(0x22, 0xC5, 0x5E)   # green
    if lr < 0.65: return RGBColor(0x6B, 0xAD, 0x35)   # yellow-green
    if lr < 0.75: return RGBColor(0xCC, 0xAA, 0x22)   # yellow
    if lr < 0.85: return RGBColor(0xE0, 0x7A, 0x22)   # orange
    if lr < 1.0: return RGBColor(0xE0, 0x44, 0x22)    # red-orange
    return RGBColor(0xEF, 0x44, 0x44)                  # red

def lr_bg(lr):
    if lr < 0.5: return RGBColor(0x15, 0x3D, 0x20)
    if lr < 0.65: return RGBColor(0x25, 0x3D, 0x15)
    if lr < 0.75: return RGBColor(0x3D, 0x35, 0x15)
    if lr < 0.85: return RGBColor(0x3D, 0x28, 0x15)
    if lr < 1.0: return RGBColor(0x3D, 0x1A, 0x15)
    return RGBColor(0x3D, 0x15, 0x15)

for i in range(rows):
    for j in range(cols):
        cell = table.cell(i, j)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.name = "Arial"
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.color.rgb = TEXT_WHITE

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        if i == 0:  # header row
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_SURFACE2
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.color.rgb = TEXT_MUTED
        elif j == 0:  # row labels
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_SURFACE
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
        else:
            lr = heatmap[AGE_GROUPS[i-1]][CAR_ERAS[j-1]]["lr"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = lr_bg(lr)
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = lr_color(lr)
                p.font.bold = True


# ── SLIDE 6: Gender Comparison ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PINK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Gender Analysis", font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "Loss ratio comparison between male and female policyholders by age", font_size=14, color=TEXT_MUTED)

# Gender comparison chart
chart_data_g = CategoryChartData()
chart_data_g.categories = AGE_GROUPS

male_lr = []
female_lr = []
for ag in AGE_GROUPS:
    m_rows = [d for d in data if age_group(d["age"]) == ag and d["gender"] == "Male"]
    f_rows = [d for d in data if age_group(d["age"]) == ag and d["gender"] == "Female"]
    m_prem = sum(r["premium"] for r in m_rows)
    m_loss = sum(r["loss"] for r in m_rows)
    f_prem = sum(r["premium"] for r in f_rows)
    f_loss = sum(r["loss"] for r in f_rows)
    male_lr.append(m_loss / m_prem * 100 if m_prem else 0)
    female_lr.append(f_loss / f_prem * 100 if f_prem else 0)

chart_data_g.add_series("Male", male_lr)
chart_data_g.add_series("Female", female_lr)

cf_g = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6),
    Inches(8), Inches(5.4), chart_data_g
)
cg = cf_g.chart
cg.has_legend = True
cg.legend.position = XL_LEGEND_POSITION.BOTTOM
cg.legend.font.size = Pt(12)
cg.legend.font.color.rgb = TEXT_MUTED
cg.legend.include_in_layout = False
plot_g = cg.plots[0]
plot_g.gap_width = 60
plot_g.series[0].format.fill.solid()
plot_g.series[0].format.fill.fore_color.rgb = ACCENT
plot_g.series[1].format.fill.solid()
plot_g.series[1].format.fill.fore_color.rgb = PINK
cg.value_axis.minimum_scale = 0
cg.value_axis.maximum_scale = 120
cg.value_axis.has_title = True
cg.value_axis.axis_title.text_frame.paragraphs[0].text = "Loss Ratio (%)"
cg.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
cg.value_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED
cg.value_axis.tick_labels.font.size = Pt(10)
cg.value_axis.tick_labels.font.color.rgb = TEXT_MUTED
cg.category_axis.tick_labels.font.size = Pt(11)
cg.category_axis.tick_labels.font.color.rgb = TEXT_MUTED

# Right side insight
card_x = Inches(9)
card_w2 = Inches(3.8)
insights_g = [
    ("Male Drivers", f"{by_gender['Male']['count']:,} policies", f"LR: {by_gender['Male']['loss_ratio']*100:.0f}%", ACCENT),
    ("Female Drivers", f"{by_gender['Female']['count']:,} policies", f"LR: {by_gender['Female']['loss_ratio']*100:.0f}%", PINK),
]
for i, (title, sub1, sub2, color) in enumerate(insights_g):
    y = Inches(2.0) + i * Inches(2.2)
    add_shape(slide, card_x, y, card_w2, Inches(1.8), BG_SURFACE, BORDER)
    add_rect(slide, card_x, y, Inches(0.06), Inches(1.8), color)
    add_text_box(slide, card_x + Inches(0.3), y + Inches(0.2), Inches(3.3), Inches(0.3),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, card_x + Inches(0.3), y + Inches(0.65), Inches(3.3), Inches(0.3),
                 sub1, font_size=13, color=TEXT_WHITE)
    add_text_box(slide, card_x + Inches(0.3), y + Inches(1.1), Inches(3.3), Inches(0.3),
                 sub2, font_size=22, color=color, bold=True)


# ── SLIDE 7: AI Premium Predictor Showcase ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), GREEN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "AI-Powered Premium & Risk Predictor", font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "Interactive what-if analysis for underwriting decisions", font_size=14, color=TEXT_MUTED)

# Example scenario cards
scenarios = [
    ("Scenario A: High Risk", "22-year-old Male, 2008 car", "$1,763", "HIGH", "44%", RED),
    ("Scenario B: Low Risk", "40-year-old Female, 2023 car", "$1,200", "LOW", "28%", GREEN),
    ("Scenario C: Medium Risk", "55-year-old Male, 2015 car", "$1,433", "MEDIUM", "33%", AMBER),
]

for i, (title, profile, premium, risk, prob, color) in enumerate(scenarios):
    x = Inches(0.8) + i * Inches(4.1)
    y = Inches(1.7)
    w = Inches(3.8)
    h = Inches(3.2)

    add_shape(slide, x, y, w, h, BG_SURFACE, BORDER)
    add_rect(slide, x, y, w, Inches(0.06), color)

    add_text_box(slide, x + Inches(0.25), y + Inches(0.25), w - Inches(0.5), Inches(0.3),
                 title, font_size=14, color=color, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.65), w - Inches(0.5), Inches(0.3),
                 profile, font_size=11, color=TEXT_MUTED)

    # Premium
    add_text_box(slide, x + Inches(0.25), y + Inches(1.15), w - Inches(0.5), Inches(0.25),
                 "PREDICTED PREMIUM", font_size=9, color=TEXT_MUTED)
    add_text_box(slide, x + Inches(0.25), y + Inches(1.4), w - Inches(0.5), Inches(0.4),
                 premium, font_size=26, color=TEXT_WHITE, bold=True)

    # Risk badge
    add_text_box(slide, x + Inches(0.25), y + Inches(2.0), Inches(1.5), Inches(0.25),
                 "RISK LEVEL", font_size=9, color=TEXT_MUTED)
    badge = add_shape(slide, x + Inches(0.25), y + Inches(2.3), Inches(1.2), Inches(0.4), color)
    add_text_box(slide, x + Inches(0.25), y + Inches(2.3), Inches(1.2), Inches(0.4),
                 risk, font_size=12, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Claim probability
    add_text_box(slide, x + Inches(1.8), y + Inches(2.0), Inches(1.8), Inches(0.25),
                 "CLAIM PROBABILITY", font_size=9, color=TEXT_MUTED)
    add_text_box(slide, x + Inches(1.8), y + Inches(2.3), Inches(1.8), Inches(0.4),
                 prob, font_size=22, color=color, bold=True)

# Bottom insight box
add_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8), BG_SURFACE, BORDER)
add_text_box(slide, Inches(1.1), Inches(5.4), Inches(5), Inches(0.3),
             "How It Works", font_size=16, color=GREEN, bold=True)

steps = [
    ("1.", " Underwriter inputs customer profile (age, gender, car model year)"),
    ("2.", " AI model calculates predicted premium using actuarial risk factors"),
    ("3.", " Risk level and claim probability are computed in real-time"),
    ("4.", " Results compared against portfolio averages for context"),
]
for i, (num, desc) in enumerate(steps):
    y_pos = Inches(5.85) + i * Inches(0.28)
    txbox = slide.shapes.add_textbox(Inches(1.1), y_pos, Inches(11), Inches(0.28))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = num
    r1.font.size = Pt(11)
    r1.font.color.rgb = GREEN
    r1.font.bold = True
    r1.font.name = "Arial"
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(11)
    r2.font.color.rgb = TEXT_WHITE
    r2.font.name = "Arial"


# ── SLIDE 8: Fraud Detection ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "AI-Powered Fraud & Anomaly Detection", font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "Automated flagging of suspicious claims using score-based rules", font_size=14, color=TEXT_MUTED)

# Fraud KPI cards
fraud_kpis = [
    ("3,539", "Claims Analyzed", CYAN),
    ("95", "Flagged Suspicious", RED),
    ("2.7%", "Flag Rate", AMBER),
    ("5", "Detection Rules", ACCENT_LIGHT),
    ("8", "Highest Risk Score", RED),
]

for i, (value, label, color) in enumerate(fraud_kpis):
    x = Inches(0.8) + i * (Inches(2.2) + Inches(0.27))
    y = Inches(1.6)
    add_shape(slide, x, y, Inches(2.2), Inches(1.3), BG_SURFACE, BORDER)
    add_rect(slide, x + Inches(0.15), y + Inches(0.12), Inches(0.35), Inches(0.04), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.3), Inches(1.9), Inches(0.25),
                 label, font_size=10, color=TEXT_MUTED)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(1.9), Inches(0.5),
                 value, font_size=26, color=color, bold=True)

# Detection rules table
rules_data = [
    ("Extreme Loss Ratio", "> 15x premium", "+3", "18", RED),
    ("Statistical Outlier", "> mean + 3\u03c3 by age group", "+2", "59", AMBER),
    ("New Car, High Loss", "Car \u2265 2022, loss > $10K", "+2", "28", AMBER),
    ("Young Driver Extreme", "Age < 25, loss > $15K", "+2", "8", AMBER),
    ("Premium-Loss Mismatch", "Top 5% loss, bottom 25% premium", "+1", "39", RGBColor(0xCC, 0xAA, 0x22)),
]

rules_table = slide.shapes.add_table(6, 4, Inches(0.8), Inches(3.2), Inches(7.5), Inches(3.5))
rtbl = rules_table.table
rtbl.columns[0].width = Inches(2.5)
rtbl.columns[1].width = Inches(2.7)
rtbl.columns[2].width = Inches(1.0)
rtbl.columns[3].width = Inches(1.3)

headers = ["Detection Rule", "Criteria", "Score", "Flagged"]
for j, h in enumerate(headers):
    cell = rtbl.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = BG_SURFACE2
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER if j >= 2 else PP_ALIGN.LEFT
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

for i, (name, criteria, score, count, color) in enumerate(rules_data):
    for j, val in enumerate([name, criteria, score, count]):
        cell = rtbl.cell(i + 1, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_SURFACE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = "Arial"
            p.font.color.rgb = color if j == 0 else TEXT_WHITE
            p.font.bold = (j == 0)
            p.alignment = PP_ALIGN.CENTER if j >= 2 else PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# Right side: top flagged case
card_x = Inches(8.8)
add_shape(slide, card_x, Inches(3.2), Inches(4), Inches(3.5), BG_SURFACE, BORDER)
add_rect(slide, card_x, Inches(3.2), Inches(4), Inches(0.06), RED)
add_text_box(slide, card_x + Inches(0.25), Inches(3.45), Inches(3.5), Inches(0.3),
             "Highest Risk Case", font_size=14, color=RED, bold=True)
add_text_box(slide, card_x + Inches(0.25), Inches(3.85), Inches(3.5), Inches(0.25),
             "Customer #9205  |  Risk Score: 8/8", font_size=11, color=TEXT_MUTED)

flagged_details = [
    ("Profile:", "54-year-old Female, 2025 car"),
    ("Premium:", "$996/year"),
    ("Claim:", "$21,147 (21.2x loss ratio)"),
    ("Flags:", "4 of 5 rules triggered"),
]
for i, (label, val) in enumerate(flagged_details):
    y_pos = Inches(4.3) + i * Inches(0.42)
    txbox = slide.shapes.add_textbox(card_x + Inches(0.25), y_pos, Inches(3.5), Inches(0.35))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = label + " "
    r1.font.size = Pt(11)
    r1.font.color.rgb = TEXT_MUTED
    r1.font.bold = True
    r1.font.name = "Arial"
    r2 = p.add_run()
    r2.text = val
    r2.font.size = Pt(11)
    r2.font.color.rgb = TEXT_WHITE
    r2.font.name = "Arial"

# Flag distribution chart
chart_data_f = CategoryChartData()
chart_data_f.categories = ["Score 1-2", "Score 3-4", "Score 5-6", "Score 7-8"]
chart_data_f.add_series("Customers", [53, 19, 21, 2])

cf_fraud = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, card_x, Inches(5.95), Inches(4), Inches(1.7), chart_data_f
)
c_fraud = cf_fraud.chart
c_fraud.has_legend = False
plot_fraud = c_fraud.plots[0]
plot_fraud.gap_width = 60
score_colors = [GREEN, AMBER, ORANGE, RED]
for i, sc in enumerate(score_colors):
    plot_fraud.series[0].points[i].format.fill.solid()
    plot_fraud.series[0].points[i].format.fill.fore_color.rgb = sc
c_fraud.value_axis.tick_labels.font.size = Pt(9)
c_fraud.value_axis.tick_labels.font.color.rgb = TEXT_MUTED
c_fraud.category_axis.tick_labels.font.size = Pt(9)
c_fraud.category_axis.tick_labels.font.color.rgb = TEXT_MUTED


# ── SLIDE 9: NL Data Chatbot ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Natural Language Data Q&A", font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "Ask questions about your portfolio in plain English — powered by Claude AI", font_size=14, color=TEXT_MUTED)

# Chat mockup - left side
chat_x = Inches(0.8)
chat_w = Inches(6.5)
add_shape(slide, chat_x, Inches(1.7), chat_w, Inches(5.3), BG_SURFACE, BORDER)
add_text_box(slide, chat_x + Inches(0.15), Inches(1.78), Inches(3), Inches(0.3),
             "Insurance Data Q&A", font_size=12, color=CYAN, bold=True)

# Simulated chat messages
chat_messages = [
    ("user", "Which age group has the highest loss ratio?"),
    ("bot", "The 65+ age group has the highest loss ratio at 85.6%, followed by 18-24 at 74.8%. The lowest is 25-34 at 65.5%."),
    ("user", "How many customers were flagged for fraud?"),
    ("bot", "95 customers were flagged across 5 detection rules, representing 2.7% of all claims. The highest risk score was 8/8."),
    ("user", "What premium should I charge a 22-year-old male with a 2008 car?"),
    ("bot", "Based on portfolio risk factors: predicted premium $1,763/yr, HIGH risk level, 44.3% claim probability — 18.7% above average."),
]

y_chat = Inches(2.2)
for sender, text in chat_messages:
    if sender == "user":
        # User message - right aligned, accent bg
        msg_w = Inches(4.5)
        msg_x = chat_x + chat_w - msg_w - Inches(0.2)
        msg_shape = add_shape(slide, msg_x, y_chat, msg_w, Inches(0.42), ACCENT)
        add_text_box(slide, msg_x + Inches(0.1), y_chat + Inches(0.05), msg_w - Inches(0.2), Inches(0.35),
                     text, font_size=9, color=TEXT_WHITE)
        y_chat += Inches(0.52)
    else:
        # Bot message - left aligned, surface2 bg
        msg_w = Inches(5.2)
        msg_x = chat_x + Inches(0.2)
        lines = len(text) / 70  # rough line count
        msg_h = Inches(0.3) + Inches(0.18) * max(lines, 1)
        msg_shape = add_shape(slide, msg_x, y_chat, msg_w, msg_h, BG_SURFACE2, BORDER)
        add_text_box(slide, msg_x + Inches(0.1), y_chat + Inches(0.03), msg_w - Inches(0.2), msg_h - Inches(0.06),
                     text, font_size=9, color=TEXT_WHITE)
        y_chat += msg_h + Inches(0.1)

# Right side: how it works
info_x = Inches(7.8)
info_w = Inches(5)
add_shape(slide, info_x, Inches(1.7), info_w, Inches(2.5), BG_SURFACE, BORDER)
add_text_box(slide, info_x + Inches(0.25), Inches(1.9), Inches(4.5), Inches(0.3),
             "How It Works", font_size=16, color=CYAN, bold=True)

how_steps = [
    ("1.", " User asks a question in plain English"),
    ("2.", " Data context (10K records, pre-aggregated) sent to Claude"),
    ("3.", " AI analyzes and responds with precise numbers"),
    ("4.", " No SQL, no coding — just conversation"),
]
for i, (num, desc) in enumerate(how_steps):
    y_pos = Inches(2.35) + i * Inches(0.38)
    txbox = slide.shapes.add_textbox(info_x + Inches(0.25), y_pos, Inches(4.5), Inches(0.35))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = num
    r1.font.size = Pt(11)
    r1.font.color.rgb = CYAN
    r1.font.bold = True
    r1.font.name = "Arial"
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(11)
    r2.font.color.rgb = TEXT_WHITE
    r2.font.name = "Arial"

# Example questions box
add_shape(slide, info_x, Inches(4.5), info_w, Inches(2.5), BG_SURFACE, BORDER)
add_text_box(slide, info_x + Inches(0.25), Inches(4.65), Inches(4.5), Inches(0.3),
             "Example Questions You Can Ask", font_size=14, color=AMBER, bold=True)

example_qs = [
    '"What is the average premium for drivers under 25?"',
    '"Compare loss ratios between male and female drivers"',
    '"Which car era has the most claims?"',
    '"Show the top 5 largest claims in the portfolio"',
    '"Is our overall loss ratio healthy?"',
]
for i, q in enumerate(example_qs):
    y_pos = Inches(5.1) + i * Inches(0.35)
    add_text_box(slide, info_x + Inches(0.25), y_pos, Inches(4.5), Inches(0.3),
                 q, font_size=10, color=TEXT_MUTED)


# ── SLIDE 10: PDF Report Showcase ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_LIGHT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Automated PDF Report Generation", font_size=32, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "AI generates a full actuarial-quality report from raw data in minutes", font_size=14, color=TEXT_MUTED)

# Report pages mockup
pages = [
    ("Page 1", "Cover Page", "Professional title page with\nbranding and date", ACCENT_LIGHT),
    ("Page 2", "Executive Summary", "KPI table, narrative insights,\nkey findings bullets", CYAN),
    ("Page 3", "Age Analysis", "Loss ratio & claims frequency\ncharts with commentary", AMBER),
    ("Page 4", "Risk Segmentation", "Heatmap matrix + car era\ntrend analysis", RED),
    ("Page 5", "Gender & Premium", "Grouped bar charts +\ndistribution histogram", PINK),
    ("Page 6", "Recommendations", "6 actionable strategies with\nprojected impact table", GREEN),
]

for i, (pg_num, title, desc, color) in enumerate(pages):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(1.7) + row * Inches(2.7)
    w = Inches(3.8)
    h = Inches(2.4)

    add_shape(slide, x, y, w, h, BG_SURFACE, BORDER)
    add_rect(slide, x, y, w, Inches(0.06), color)

    # Page number badge
    badge = add_shape(slide, x + Inches(0.2), y + Inches(0.25), Inches(0.7), Inches(0.35), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.25), Inches(0.7), Inches(0.35),
                 pg_num, font_size=10, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), w - Inches(0.4), Inches(0.3),
                 title, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), w - Inches(0.4), Inches(0.9),
                 desc, font_size=11, color=TEXT_MUTED)


# ── SLIDE 11: What AI Agents Built ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.5),
             "What AI Agents Can Do For Us", font_size=28, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.35),
             "Everything in this demo was generated by Claude Code in a single session", font_size=12, color=TEXT_MUTED)

capabilities = [
    ("Mock Data", "10K realistic\nrecords", ACCENT_LIGHT),
    ("Dashboard", "Charts, heatmap\n& filters", CYAN),
    ("Predictor", "What-if premium\ncalculator", GREEN),
    ("Presentation", "Auto-generated\nslides", AMBER),
    ("PDF Report", "6-page actuarial\nreport", ACCENT_LIGHT),
    ("Fraud Detection", "Score-based\nanomalies", ORANGE),
    ("Data Chatbot", "NL Q&A\nvia AI", PINK),
]

for i, (title, desc, color) in enumerate(capabilities):
    x = Inches(0.5) + i * Inches(1.82)
    y = Inches(1.25)
    w = Inches(1.65)
    h = Inches(2.0)

    add_shape(slide, x, y, w, h, BG_SURFACE, BORDER)
    add_rect(slide, x, y, w, Inches(0.05), color)

    # Number badge
    add_shape(slide, x + Inches(0.12), y + Inches(0.13), Inches(0.35), Inches(0.35), color)
    add_text_box(slide, x + Inches(0.12), y + Inches(0.13), Inches(0.35), Inches(0.35),
                 str(i + 1), font_size=11, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.1), y + Inches(0.6), w - Inches(0.2), Inches(0.28),
                 title, font_size=11, color=color, bold=True)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.9), w - Inches(0.2), Inches(0.8),
                 desc, font_size=9, color=TEXT_MUTED)

# ── Speed Comparison Table ──
speed_label_y = Inches(3.45)
add_text_box(slide, Inches(0.5), speed_label_y, Inches(12), Inches(0.3),
             "AI Speed Comparison  —  Minutes vs Days", font_size=13, color=ACCENT_LIGHT, bold=True)

speed_data = [
    ("Deliverable",          "Traditional Approach",           "With AI Agent"),
    ("Mock Data Generation", "1–2 days  (data team setup)",    "~5 min"),
    ("Interactive Dashboard","1–2 weeks  (frontend dev)",      "~30 min"),
    ("Premium Predictor",    "3–5 days  (model dev)",          "~15 min"),
    ("Presentation",         "1–2 days  (manual slides)",      "~10 min"),
    ("PDF Actuarial Report", "3–5 days  (report team)",        "~20 min"),
    ("Fraud Detection",      "2–4 weeks  (data science team)", "~30 min"),
    ("NL Data Chatbot",      "1–2 months  (dev team)",         "~45 min"),
]

tbl_shape = slide.shapes.add_table(
    len(speed_data), 3,
    Inches(0.5), Inches(3.82),
    Inches(12.33), Inches(3.5)
)
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.8)
tbl.columns[1].width = Inches(6.1)
tbl.columns[2].width = Inches(3.43)

GREEN_BG  = RGBColor(0x0D, 0x22, 0x14)
RED_BG    = RGBColor(0x22, 0x10, 0x10)
ORANGE_FG = RGBColor(0xF9, 0x73, 0x16)

for r, (task, trad, ai_t) in enumerate(speed_data):
    is_header = (r == 0)
    for c, text in enumerate([task, trad, ai_t]):
        cell = tbl.cell(r, c)
        tf = cell.text_frame
        tf.word_wrap = True
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
        run = para.add_run()
        run.text = text
        run.font.size = Pt(10)
        run.font.bold = is_header or (c == 0)

        if is_header:
            run.font.color.rgb = TEXT_WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_SURFACE2
        elif c == 2:                      # AI time — green
            run.font.color.rgb = GREEN
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREEN_BG
        elif c == 1:                      # Traditional — orange/muted
            run.font.color.rgb = ORANGE_FG
            cell.fill.solid()
            cell.fill.fore_color.rgb = RED_BG
        else:                             # Task name
            run.font.color.rgb = TEXT_WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_SURFACE


# ── SLIDE 9: Thank You ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), CYAN)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.0),
             "Thank You", font_size=48, color=ACCENT_LIGHT, bold=True,
             alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.02), BORDER)

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "Questions & Discussion", font_size=20, color=TEXT_MUTED,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
             "All materials available: dashboard, data, reports, fraud detection, and chatbot",
             font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "Powered by Claude Code  |  Anthropic", font_size=12, color=ACCENT,
             alignment=PP_ALIGN.CENTER)


# ── Save ──
output_path = "Auto_Insurance_AI_Demo_V2.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
